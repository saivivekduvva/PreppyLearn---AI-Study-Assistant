import os
import re
import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
from app.utils.logger import logger
from app.utils.exceptions import FileNotFoundException, DocumentExtractionException, InvalidDocumentException
from pdfminer.pdfparser import PDFSyntaxError

class ExtractionService:
    def __init__(self):
        self.upload_dir = os.path.join(os.getcwd(), "uploads")

    def extract_text(self, filename: str) -> str:
        """
        Extracts clean text from a file located in the uploads directory.
        Handles PDFs, DOCX, PPTX, and TXT files.
        """
        file_path = os.path.join(self.upload_dir, filename)
        
        if not os.path.exists(file_path):
            logger.warning(f"Extraction failed: File not found -> {file_path}")
            raise FileNotFoundException(detail=f"File '{filename}' not found.")
            
        try:
            logger.info(f"Starting text extraction from {filename}...")
            
            ext = os.path.splitext(filename)[1].lower()
            if ext == '.pdf':
                raw_text = self._extract_pdf(file_path)
            elif ext in ['.docx', '.doc']:
                raw_text = self._extract_docx(file_path)
            elif ext in ['.pptx', '.ppt']:
                raw_text = self._extract_pptx(file_path)
            elif ext == '.txt':
                raw_text = self._extract_txt(file_path)
            else:
                raise Exception(f"Unsupported file type: {ext}")
            
            if not raw_text:
                logger.warning(f"No text could be extracted from {filename}.")
                return ""
                
            # Clean up whitespace to ensure high quality LLM context
            # Replace multiple spaces with a single space
            clean_text = re.sub(r'[ \t]+', ' ', raw_text)
            # Replace 3 or more newlines with 2 newlines (to preserve paragraphs but remove huge vertical gaps)
            clean_text = re.sub(r'\n{3,}', '\n\n', clean_text)
            
            clean_text = clean_text.strip()
            
            logger.info(f"Successfully extracted {len(clean_text)} characters from {filename}.")
            return clean_text
            
        except PDFSyntaxError as e:
            logger.error(f"Invalid or corrupted PDF file {filename}: {str(e)}")
            raise InvalidDocumentException(detail="The uploaded PDF is invalid or corrupted.")
        except Exception as e:
            logger.error(f"Failed to extract text from {filename}: {str(e)}")
            raise DocumentExtractionException(detail=f"Failed to extract text: {str(e)}")

    def _extract_pdf(self, file_path: str) -> str:
        text_content = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_content.append(page_text)
        return "\n".join(text_content)

    def _extract_docx(self, file_path: str) -> str:
        doc = DocxDocument(file_path)
        text_content = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_content.append(para.text.strip())
        return "\n".join(text_content)

    def _extract_pptx(self, file_path: str) -> str:
        prs = Presentation(file_path)
        text_content = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text.strip())
        return "\n".join(text_content)

    def _extract_txt(self, file_path: str) -> str:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
